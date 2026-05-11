"""Multi-format document text extraction (PDF/DOCX/PPTX/XLSX/TXT/MD) with optional OCR."""
from __future__ import annotations

import os

from pypdf import PdfReader

from eigenmind.config import ocr_available

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import pytesseract
    from pdf2image import convert_from_path
except ImportError:
    pytesseract = None
    convert_from_path = None


def _extract_pdf(filepath: str, log) -> str:
    """Extract text from a PDF, falling back to OCR if the text layer looks like a scan."""
    reader = PdfReader(filepath)
    extracted = ""
    for page in reader.pages:
        try:
            extracted += page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            log(f"  -> Warning: Could not extract text from a page: {e}")

    n_pages = len(reader.pages)
    looks_scanned = n_pages > 0 and len(extracted.strip()) < (20 * n_pages)
    if not looks_scanned:
        return extracted

    if not ocr_available() or pytesseract is None or convert_from_path is None:
        log("  -> Warning: PDF seems to be a scan, but OCR libs (pdf2image, pytesseract) "
            "are missing or TESSDATA_PREFIX is not set.")
        return extracted

    try:
        log(f"  -> Detected potential scan. Running OCR on {os.path.basename(filepath)}...")
        ocr_text = ""
        for img in convert_from_path(filepath):
            ocr_text += pytesseract.image_to_string(img) + "\n"
        log(f"  -> OCR complete. Extracted {len(ocr_text)} characters.")
        return ocr_text if len(ocr_text.strip()) > len(extracted.strip()) else extracted
    except Exception as e:  # noqa: BLE001
        log(f"  -> OCR failed (check Poppler/Tesseract PATH): {e}")
        return extracted


def _extract_docx(filepath: str, log) -> str:
    if docx is None:
        log("  -> Error: 'python-docx' not installed. Skipping .docx file.")
        return ""
    return "\n".join(p.text for p in docx.Document(filepath).paragraphs)


def _extract_pptx(filepath: str, log) -> str:
    if Presentation is None:
        log("  -> Error: 'python-pptx' not installed. Skipping .pptx file.")
        return ""
    out = ""
    for slide in Presentation(filepath).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                out += shape.text + "\n"
    return out


def _extract_xlsx(filepath: str, log) -> str:
    if pd is None:
        log("  -> Error: 'pandas'/'openpyxl' not installed. Skipping .xlsx file.")
        return ""
    out = ""
    for sheet_name, df in pd.read_excel(filepath, sheet_name=None).items():
        out += f"\n--- Sheet: {sheet_name} ---\n" + df.to_string()
    return out


def _extract_text_file(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text(filepath: str, log=lambda _msg: None) -> str:
    """Dispatch text extraction by file extension. Returns "" for unsupported types."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return _extract_pdf(filepath, log)
    if ext == ".docx":
        return _extract_docx(filepath, log)
    if ext == ".pptx":
        return _extract_pptx(filepath, log)
    if ext == ".xlsx":
        return _extract_xlsx(filepath, log)
    if ext in (".txt", ".md"):
        return _extract_text_file(filepath)
    return ""
